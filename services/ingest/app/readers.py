"""Document readers for the ingest microservice.

Format-aware parsing utilities for diverse document types with optional
multimodal OCR/captioning via Gemini. This module:
- Dispatches by file extension to dedicated parsers: PDF (pdfplumber with
    pages and table extraction), DOCX (python-docx), PPTX (python-pptx),
    HTML (BeautifulSoup+html5lib), Markdown (markdown-it-py), CSV (pandas),
    TXT (decode), and common image types.
- For images and PDF pages when enabled, uses Gemini 2.5 Pro/Flash as a
    cloud vision endpoint for OCR/captioning and merges the results into the
    returned text to create a unified pseudo-text representation.
- Enforces soft limits (file size/pages) via Settings; large files can be
    truncated progressively.

The top-level `parse_document(content, filename) -> str` returns a best-effort
normalized text suitable for downstream chunking. Metadata such as page/section
will be attached later during chunk construction in the caller.
"""

from __future__ import annotations

import io
import os
from typing import Generator, Tuple

import pandas as pd
import pdfplumber
from bs4 import BeautifulSoup
from PIL import Image

from shared.settings import Settings

_settings = Settings()
_GEMINI_KEY = _settings.gemini_api_key
_GEMINI_ENABLE_MM = _settings.gemini_multimodal_enabled and bool(_GEMINI_KEY)
_GEMINI_FAST = _settings.gemini_fast_model
_GEMINI_REASON = _settings.gemini_reasoning_model

genai = None
if _GEMINI_ENABLE_MM and not _settings.offline_mode:
    try:
        import google.generativeai as genai  # type: ignore

        genai.configure(api_key=_GEMINI_KEY)
    except Exception:
        genai = None


def parse_document(content: bytes, filename: str) -> str:
    """Parse a document into normalized text.

    Dispatches based on file extension; falls back to robust text decode.
    For PDFs, uses pdfplumber with page iteration, basic layout handling,
    and table extraction. Optionally, renders page images for OCR/captioning
    via Gemini to augment text.
    """
    name = filename.lower()
    ext = os.path.splitext(name)[1]
    max_mb = _settings.ingest_max_file_mb
    if max_mb and len(content) > max_mb * 1024 * 1024:
        content = content[: max_mb * 1024 * 1024]

    if ext in {".pdf"}:
        return _parse_pdf(content)
    if ext in {".docx"}:
        return _parse_docx(content)
    if ext in {".pptx"}:
        return _parse_pptx(content)
    if ext in {".html", ".htm"}:
        return _parse_html(content)
    if ext in {".md", ".markdown"}:
        return _parse_markdown(content)
    if ext in {".csv"}:
        return _parse_csv(content)
    if ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"}:
        return _parse_image(content)
    try:
        return content.decode("utf-8", errors="replace")
    except Exception:
        return content.decode("latin-1", errors="replace")


def _table_to_csv(tbl: list[list[object]]) -> str:
    """Convert a pdfplumber table (list of rows) into CSV text.

    Falls back gracefully on non-string cells. No quoting to keep simple; cells
    containing commas are left as-is which is acceptable for downstream LLMs.
    """
    lines: list[str] = []
    for row in tbl:
        if not row:
            continue
        cells = ["" if c is None else str(c).replace("\n", " ") for c in row]
        lines.append(",".join(cells))
    return "\n".join(lines)


def _parse_pdf(content: bytes) -> str:
    output_lines: list[str] = []
    max_pages = max(1, _settings.ingest_max_pages)
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page_idx, page in enumerate(pdf.pages[:max_pages]):
            try:
                text = page.extract_text(x_tolerance=1.5, y_tolerance=1.5) or ""
            except Exception:
                text = ""
            tables_text = []
            try:
                tables = page.extract_tables()
                for ti, tbl in enumerate(tables or []):
                    if not tbl:
                        continue
                    table_id = f"p{page_idx + 1}_t{ti + 1}"
                    if _settings.normalize_tables_to_csv:
                        csv_text = _table_to_csv(tbl)
                    else:
                        csv_lines = []
                        for row in tbl:
                            if row:
                                csv_lines.append(
                                    "\t".join(str(c) if c else "" for c in row)
                                )
                        csv_text = "\n".join(csv_lines)
                    tables_text.append(f"[table id={table_id}]\n{csv_text}\n[/table]")
            except Exception:
                pass
            header = f"\n\n=== Page {page_idx + 1} ===\n"
            block = text.strip()
            if tables_text:
                block += "\n" + "\n".join(tables_text)
            if _settings.pdf_render_images and _GEMINI_ENABLE_MM and genai is not None:
                try:
                    pil = page.to_image(resolution=150).original
                    ocr_txt = _gemini_image_to_text(pil)
                    if ocr_txt:
                        block += f"\n[vision_ocr]\n{ocr_txt.strip()}"
                except Exception:
                    pass
            if block:
                output_lines.append(header + block)
    return "\n".join(output_lines).strip()


def stream_pdf_pages(
    content: bytes, max_pages: int | None = None
) -> Generator[Tuple[int, str], None, None]:
    """Yield (page_number, page_text) for a PDF, with table markers.

    Designed for NDJSON streaming in the ingest service to handle very large
    PDFs without materializing the entire parsed text in memory.
    """
    limit = max_pages or _settings.ingest_max_pages
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for page_idx, page in enumerate(pdf.pages[: max(1, limit)]):
            try:
                text = page.extract_text(x_tolerance=1.5, y_tolerance=1.5) or ""
            except Exception:
                text = ""
            tables_text = []
            try:
                tables = page.extract_tables()
                for ti, tbl in enumerate(tables or []):
                    if not tbl:
                        continue
                    table_id = f"p{page_idx + 1}_t{ti + 1}"
                    if _settings.normalize_tables_to_csv:
                        csv_text = _table_to_csv(tbl)
                    else:
                        csv_lines = []
                        for row in tbl:
                            if row:
                                csv_lines.append(
                                    "\t".join(str(c) if c else "" for c in row)
                                )
                        csv_text = "\n".join(csv_lines)
                    tables_text.append(f"[table id={table_id}]\n{csv_text}\n[/table]")
            except Exception:
                pass
            header = f"\n\n=== Page {page_idx + 1} ===\n"
            block = text.strip()
            if tables_text:
                block += "\n" + "\n".join(tables_text)
            if _settings.pdf_render_images and _GEMINI_ENABLE_MM and genai is not None:
                try:
                    pil = page.to_image(resolution=150).original
                    ocr_txt = _gemini_image_to_text(pil)
                    if ocr_txt:
                        block += f"\n[vision_ocr]\n{ocr_txt.strip()}"
                except Exception:
                    pass
            if block:
                yield (page_idx + 1, header + block)


def _parse_docx(content: bytes) -> str:
    try:
        from docx import Document  # type: ignore

        doc = Document(io.BytesIO(content))
        paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        for tbl in doc.tables:
            for row in tbl.rows:
                paras.append("\t".join(cell.text.strip() for cell in row.cells))
        return "\n".join(paras)
    except Exception:
        return content.decode("utf-8", errors="replace")


def _parse_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(content))
        slides_out: list[str] = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shp in slide.shapes:
                if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                    for p in shp.text_frame.paragraphs:
                        txt = "".join(run.text for run in p.runs).strip()
                        if txt:
                            texts.append(txt)
                if hasattr(shp, "image") and _GEMINI_ENABLE_MM and genai is not None:
                    try:
                        img_bytes = shp.image.blob
                        pil = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                        cap = _gemini_image_to_text(pil)
                        if cap:
                            texts.append(f"[image] {cap}")
                    except Exception:
                        pass
            if texts:
                slides_out.append(f"\n\n=== Slide {i + 1} ===\n" + "\n".join(texts))
        return "\n".join(slides_out)
    except Exception:
        return content.decode("utf-8", errors="replace")


def _parse_html(content: bytes) -> str:
    soup = BeautifulSoup(content, "html5lib")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    lines: list[str] = []
    for el in soup.find_all(
        ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th"]
    ):
        txt = el.get_text(" ", strip=True)
        if txt:
            lines.append(txt)
    return "\n".join(lines)


def _parse_markdown(content: bytes) -> str:
    text = content.decode("utf-8", errors="replace")
    return text


def _parse_csv(content: bytes) -> str:
    try:
        df = pd.read_csv(io.BytesIO(content))
        return df.to_csv(index=False)
    except Exception:
        try:
            df = pd.read_excel(io.BytesIO(content))
            return df.to_csv(index=False)
        except Exception:
            return content.decode("utf-8", errors="replace")


def _parse_image(content: bytes) -> str:
    try:
        pil = Image.open(io.BytesIO(content)).convert("RGB")
    except Exception:
        return ""
    if _GEMINI_ENABLE_MM and genai is not None:
        try:
            return _gemini_image_to_text(pil)
        except Exception:
            pass
    return ""


def _gemini_image_to_text(pil_image: Image.Image) -> str:
    """Send an image to Gemini to obtain OCR/caption text."""
    if genai is None:
        return ""
    try:
        model_name = _GEMINI_REASON
        model = genai.GenerativeModel(model_name)
        prompt = "Describe or transcribe the image succinctly. Include visible text exactly when present."
        resp = model.generate_content([pil_image, prompt])
        return getattr(resp, "text", "") or ""
    except Exception:
        return ""
